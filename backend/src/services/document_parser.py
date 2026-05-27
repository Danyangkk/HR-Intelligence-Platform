from __future__ import annotations

import io
import re
from pathlib import Path


def extract_text(filename: str, data: bytes) -> str:
    ext = Path(filename).suffix.lower()
    if ext in {".txt", ".md"}:
        return data.decode("utf-8", errors="ignore")
    if ext == ".docx":
        return _extract_docx(data)
    if ext == ".pdf":
        return _extract_pdf(data)
    if ext in {".pptx", ".ppt"}:
        return _extract_pptx(data)
    raise ValueError(f"unsupported file type: {ext or 'unknown'}")


def _extract_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n".join(parts)


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        text = (page.extract_text() or "").strip()
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)


def chunk_text(text: str, *, chunk_size: int = 500, overlap: int = 75) -> list[str]:
    text = re.sub(r"\r\n?", "\n", text.strip())
    if not text:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_size)
        piece = text[start:end].strip()
        if piece:
            chunks.append(piece)
        if end >= len(text):
            break
        start = max(end - overlap, start + 1)
    return chunks
